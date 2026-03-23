#!/usr/bin/env python3
"""
03_rename_pages.py — Copy/rename page images to a stable naming scheme.

Reads source images from a folder (or extracts from a PPTX's embedded images),
renames them according to a stable canonical pattern, and COPIES them (never
moves or deletes originals) to early_data/preprocessed_pages/.

Target naming pattern:
  FBW_V03_N01_p001.jpeg         — Footbag World Vol. 3 No. 1, single page
  FBW_V05_N02_p001.jpeg         — Footbag World Vol. 5 No. 2, page 1
  FBW_V10_N02_p001a.jpeg        — Footbag World Vol. 10 No. 2, page 1a
  IFAB_WH_p001.jpeg             — IFAB Rulebook Worlds History, page 1

Source filename formats handled:
  "Footbag World Vol. 10 No. 2 page 1a.jpeg"
  "Footbag World Vol. 3, No. 2 Page 1.jpeg"
  "PREP_Footbag World Vol. 2 No. 1.png"
  "IFAB Rulebook Worlds History Page 1.jpeg"

Usage:
  # Copy from a source folder of raw images:
  python early_data/scripts/03_rename_pages.py --source <SOURCE_DIR>

  # Extract embedded images directly from the PPTX:
  python early_data/scripts/03_rename_pages.py --pptx <PPTX_FILE>

  # Dry run (show what would be copied, don't write anything):
  python early_data/scripts/03_rename_pages.py --source <DIR> --dry-run

Default output: early_data/preprocessed_pages/
"""

import argparse
import re
import shutil
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_PPTX = REPO_ROOT / "raw_images_for_manual_crop_rotation2.pptx"
DEFAULT_OUTPUT = REPO_ROOT / "early_data" / "preprocessed_pages"

# ── Filename parsing ───────────────────────────────────────────────────────────

FBW_PATTERN = re.compile(
    r"(?:PREP_)?Footbag World Vol[.\s]+(\d+)[,\s]+No[.\s]+(\d+)"
    r"(?:\s+[Pp]age\s+(\w+))?",
    re.IGNORECASE,
)

IFAB_PATTERN = re.compile(r"IFAB\s+Rulebook", re.IGNORECASE)
PAGE_NUM_PATTERN = re.compile(r"[Pp]age\s+(\w+)")


def parse_to_canonical(filename: str) -> str | None:
    """
    Convert a source filename to canonical stable name (without extension).
    Returns None if the filename cannot be parsed reliably.
    """
    stem = Path(filename).stem

    # IFAB Rulebook
    if IFAB_PATTERN.search(filename):
        m = PAGE_NUM_PATTERN.search(filename)
        page = m.group(1).zfill(3) if m else "001"
        return f"IFAB_WH_p{page}"

    # Footbag World
    m = FBW_PATTERN.search(filename)
    if m:
        vol = int(m.group(1))
        num = int(m.group(2))
        page_raw = m.group(3)  # "1", "1a", "1b", "2", etc. or None

        if page_raw is None:
            page_str = "p001"
        elif page_raw.isdigit():
            page_str = f"p{int(page_raw):03d}"
        else:
            # e.g. "1a", "1b" — keep suffix as-is after zero-padding numeric prefix
            num_part = re.match(r"(\d+)(\D*)", page_raw)
            if num_part:
                page_str = f"p{int(num_part.group(1)):03d}{num_part.group(2)}"
            else:
                page_str = f"p{page_raw}"

        return f"FBW_V{vol:02d}_N{num:02d}_{page_str}"

    return None  # unrecognized — skip or warn


# ── PPTX image extraction ─────────────────────────────────────────────────────

def extract_from_pptx(pptx_path: Path, output_dir: Path, dry_run: bool) -> int:
    """Extract embedded slide images from the PPTX using slide label as filename hint."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("ERROR: python-pptx not installed.")
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    copied = 0
    skipped = 0

    for i in range(len(prs.slides)):
        try:
            slide = prs.slides[i]
        except Exception as e:
            print(f"  WARNING: Slide {i+1} inaccessible: {e}")
            continue

        # Get the slide label (text box with filename)
        source_file = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                source_file = shape.text.strip()
                break

        if not source_file:
            print(f"  Slide {i+1}: no label text found — skipping")
            skipped += 1
            continue

        canonical = parse_to_canonical(source_file)
        if canonical is None:
            print(f"  Slide {i+1}: cannot parse filename '{source_file}' — skipping")
            skipped += 1
            continue

        # Find the image shape on this slide
        image_shapes = [
            s for s in slide.shapes
            if s.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
        ]

        if not image_shapes:
            print(f"  Slide {i+1}: no image shape found — skipping")
            skipped += 1
            continue

        # Take the first (and typically only) image
        img_shape = image_shapes[0]
        img_blob = img_shape.image.blob
        img_ext = img_shape.image.ext  # 'jpeg', 'png', etc.

        dest_name = f"{canonical}.{img_ext}"
        dest_path = output_dir / dest_name

        if dest_path.exists():
            print(f"  Slide {i+1}: {dest_name} already exists — skipping")
            skipped += 1
            continue

        if dry_run:
            print(f"  [DRY RUN] Slide {i+1}: {source_file!r} → {dest_name}")
        else:
            output_dir.mkdir(parents=True, exist_ok=True)
            dest_path.write_bytes(img_blob)
            print(f"  Slide {i+1}: {source_file!r} → {dest_name}")

        copied += 1

    print(f"\nDone. Extracted: {copied}, Skipped: {skipped}")
    return copied


# ── Folder copy mode ──────────────────────────────────────────────────────────

def copy_from_folder(source_dir: Path, output_dir: Path, dry_run: bool) -> int:
    """Rename and copy images from a source folder."""
    SUPPORTED = {".jpg", ".jpeg", ".png", ".tif", ".tiff"}
    images = sorted(p for p in source_dir.iterdir() if p.suffix.lower() in SUPPORTED)

    if not images:
        print(f"No images found in: {source_dir}")
        return 0

    print(f"Found {len(images)} images in: {source_dir}")
    copied = 0
    skipped = 0

    for img in images:
        canonical = parse_to_canonical(img.name)
        if canonical is None:
            print(f"  Cannot parse: {img.name!r} — skipping")
            skipped += 1
            continue

        dest_name = f"{canonical}{img.suffix.lower()}"
        dest_path = output_dir / dest_name

        if dest_path.exists():
            print(f"  {dest_name} already exists — skipping")
            skipped += 1
            continue

        if dry_run:
            print(f"  [DRY RUN] {img.name!r} → {dest_name}")
        else:
            output_dir.mkdir(parents=True, exist_ok=True)
            shutil.copy2(img, dest_path)
            print(f"  {img.name!r} → {dest_name}")

        copied += 1

    print(f"\nDone. Copied: {copied}, Skipped: {skipped}")
    return copied


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    group = parser.add_mutually_exclusive_group()
    group.add_argument("--pptx", type=Path, default=None,
                       help="Extract embedded images from PPTX (default: raw_images_for_manual_crop_rotation2.pptx)")
    group.add_argument("--source", type=Path, default=None,
                       help="Copy and rename images from a source folder")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT,
                       help=f"Output directory (default: {DEFAULT_OUTPUT})")
    parser.add_argument("--dry-run", action="store_true",
                       help="Show what would be copied without writing anything")
    args = parser.parse_args()

    if args.dry_run:
        print("DRY RUN — no files will be written\n")

    if args.source:
        if not args.source.exists():
            print(f"ERROR: Source directory not found: {args.source}")
            sys.exit(1)
        copy_from_folder(args.source, args.output, args.dry_run)
    else:
        pptx_path = args.pptx or DEFAULT_PPTX
        if not pptx_path.exists():
            print(f"ERROR: PPTX not found: {pptx_path}")
            sys.exit(1)
        extract_from_pptx(pptx_path, args.output, args.dry_run)


if __name__ == "__main__":
    main()
