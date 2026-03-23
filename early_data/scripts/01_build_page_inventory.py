#!/usr/bin/env python3
"""
01_build_page_inventory.py — Build page inventory from PowerPoint source file.

Reads a PPTX file where each slide contains:
  - one image (the cleaned page scan)
  - one text box with the original source filename as the label

Outputs:
  early_data/page_inventory/fbw_page_inventory.csv

Usage:
  python early_data/scripts/01_build_page_inventory.py [PPTX_FILE]

  Default PPTX_FILE: raw_images_for_manual_crop_rotation2.pptx
  (relative to repo root)
"""

import csv
import re
import sys
from pathlib import Path

# ── Configuration ──────────────────────────────────────────────────────────────

REPO_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_PPTX = REPO_ROOT / "raw_images_for_manual_crop_rotation2.pptx"
OUTPUT_CSV = REPO_ROOT / "early_data" / "page_inventory" / "fbw_page_inventory.csv"

# Known volume → approximate year range and primary year.
# Based on magazine_scan_index.csv cross-reference.
# No. 1 issues are typically results/worlds issues (published late in the year).
# No. 2+ issues are mid-year content.
VOL_YEAR_MAP = {
    2:  (1980, "1980–1981"),
    3:  (1981, "1981–1982"),
    4:  (1982, "1982–1983"),
    5:  (1983, "1983–1984"),
    6:  (1984, "1984–1986"),
    7:  (1985, "1985–1988"),
    8:  (1989, "1989"),
    9:  (1990, "1990"),
    10: (1991, "1991"),
    11: (1992, "1992"),
    12: (1993, "1993"),   # estimated — not in scan index
    13: (1994, "1994"),
    14: (1995, "1995"),
}

# ── Filename parsing ───────────────────────────────────────────────────────────

# Matches: "Footbag World Vol. 10 No. 2 page 1a.jpeg"
#          "Footbag World Vol. 3, No. 2 Page 1.jpeg"
#          "PREP_Footbag World Vol. 2 No. 1.png"
FBW_PATTERN = re.compile(
    r"(?:PREP_)?Footbag World Vol[.\s]+(\d+)[,\s]+No[.\s]+(\d+)"
    r"(?:\s+[Pp]age\s+(\w+))?",
    re.IGNORECASE,
)

IFAB_PATTERN = re.compile(r"IFAB\s+Rulebook", re.IGNORECASE)


def parse_filename(filename: str) -> dict:
    """Parse a slide label filename into structured fields."""
    result = {
        "issue": "",
        "page_label": "",
        "year_guess": "",
        "source_type": "",
        "notes": "",
    }

    # Strip extension for display
    stem = Path(filename).stem

    # Check for PREP_ prefix (indicates prior preprocessing pass)
    if filename.upper().startswith("PREP_"):
        result["notes"] = "PREP_ prefix: image was preprocessed before PowerPoint import"

    # IFAB Rulebook pages
    if IFAB_PATTERN.search(filename):
        m = re.search(r"[Pp]age\s+(\d+)", filename)
        page_num = m.group(1) if m else ""
        result["issue"] = "IFAB Rulebook"
        result["page_label"] = f"p{page_num}" if page_num else ""
        result["year_guess"] = ""   # unknown without further research
        result["source_type"] = "IFAB_RULEBOOK"
        return result

    # Footbag World magazine
    m = FBW_PATTERN.search(filename)
    if m:
        vol = int(m.group(1))
        num = int(m.group(2))
        page_raw = m.group(3)  # may be "1", "1a", "1b", "2", etc. or None

        result["issue"] = f"Vol. {vol} No. {num}"
        result["page_label"] = f"p{page_raw}" if page_raw else "p1"
        result["source_type"] = "FBW_IMAGE"

        if vol in VOL_YEAR_MAP:
            primary_year, year_range = VOL_YEAR_MAP[vol]
            # No. 1 issues are typically late in the year / results issue
            # No. 2+ are mid-year
            result["year_guess"] = str(primary_year) if num == 1 else year_range
        else:
            result["year_guess"] = ""
            result["notes"] = (result["notes"] + f" Vol. {vol} not in known year map.").strip()

        return result

    # Unrecognized format
    result["source_type"] = "UNKNOWN"
    result["notes"] = f"Could not parse filename: {filename}"
    return result


# ── Main ──────────────────────────────────────────────────────────────────────

def build_inventory(pptx_path: Path, output_path: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    if not pptx_path.exists():
        print(f"ERROR: PPTX file not found: {pptx_path}")
        sys.exit(1)

    print(f"Reading: {pptx_path}")
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    print(f"  {total} slides found")

    rows = []
    for i in range(total):
        try:
            slide = prs.slides[i]
        except Exception as e:
            print(f"  WARNING: Slide {i+1} inaccessible: {e}")
            rows.append({
                "slide_num": i + 1,
                "source_file": "",
                "issue": "",
                "page_label": "",
                "year_guess": "",
                "source_type": "",
                "has_results": "",
                "has_worlds_data": "",
                "notes": f"Slide inaccessible: {e}",
            })
            continue

        # Collect all non-empty text from the slide
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        # The filename label is the primary text on the slide.
        # If multiple text boxes, use the first one (filename convention holds in this deck).
        source_file = texts[0] if texts else ""

        parsed = parse_filename(source_file)

        rows.append({
            "slide_num": i + 1,
            "source_file": source_file,
            "issue": parsed["issue"],
            "page_label": parsed["page_label"],
            "year_guess": parsed["year_guess"],
            "source_type": parsed["source_type"],
            "has_results": "",      # fill during review pass
            "has_worlds_data": "",  # fill during review pass
            "notes": parsed["notes"],
        })

    # Write CSV
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fieldnames = [
        "slide_num", "source_file", "issue", "page_label",
        "year_guess", "source_type", "has_results", "has_worlds_data", "notes",
    ]
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)

    print(f"  Written: {output_path} ({len(rows)} rows)")
    unknown = [r for r in rows if r["source_type"] == "UNKNOWN"]
    if unknown:
        print(f"  WARNING: {len(unknown)} slides had unrecognized filenames:")
        for r in unknown:
            print(f"    Slide {r['slide_num']}: {r['source_file']}")


if __name__ == "__main__":
    pptx_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_PPTX
    build_inventory(pptx_path, OUTPUT_CSV)
