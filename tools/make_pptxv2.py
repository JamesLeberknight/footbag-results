import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# --- CONFIG ---
# If you are running from the /tools folder, we go up one level to find /out
ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = ROOT / "out" / "scans"
OUTPUT_PPTX = ROOT / "out" / "Footbag_Archive_Scans.pptx"

def create_pptx():
    # Create a basic presentation
    prs = Presentation()
    
    # 16:9 Widescreen dimensions (Optional, but usually looks better)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Get all jpg/jpeg files
    valid_exts = ('.jpg', '.jpeg', '.JPG', '.JPEG')
    
    if not IMG_DIR.exists():
        print(f"Error: Folder not found at {IMG_DIR}")
        return

    images = sorted([f for f in os.listdir(IMG_DIR) if f.endswith(valid_exts)])

    if not images:
        print(f"No images found in {IMG_DIR}")
        return

    print(f"Found {len(images)} images. Building PowerPoint...")

    for img_name in images:
        img_path = IMG_DIR / img_name
        
        # Add a Blank Slide (Layout 6 is standard for 'Blank')
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        try:
            # Place image slightly offset from top-left
            # We set the height to 7 inches so it fits on the 7.5 inch slide
            left = Inches(0.5)
            top = Inches(0.25)
            slide.shapes.add_picture(str(img_path), left, top, height=Inches(7.0))
            
            # Add the filename as a small text label at the bottom for reference
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.4))
            tf = txBox.text_frame
            tf.text = f"File: {img_name}"
            
        except Exception as e:
            print(f"Skipping {img_name} due to error: {e}")

    # Save to the /out folder
    prs.save(OUTPUT_PPTX)
    print(f"--- SUCCESS ---")
    print(f"Saved {len(images)} slides to: {OUTPUT_PPTX}")

if __name__ == "__main__":
    create_pptx()
