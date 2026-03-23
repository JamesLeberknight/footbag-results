import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageEnhance, ImageOps

# --- CONFIG ---
ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = ROOT / "out" / "scans"
# We save enhanced versions separately so we don't destroy your original files
ENH_DIR = ROOT / "out" / "enhanced_scans" 
OUTPUT_PPTX = ROOT / "out" / "Footbag_Archive_Enhanced_v3.pptx"

ENH_DIR.mkdir(parents=True, exist_ok=True)

def enhance_scans():
    prs = Presentation()
    # Widescreen 16:9 layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    valid_exts = ('.jpg', '.jpeg', '.JPG', '.JPEG')
    images = sorted([f for f in os.listdir(IMG_DIR) if f.endswith(valid_exts)])

    print(f"Processing {len(images)} scans...")

    for img_name in images:
        orig_path = IMG_DIR / img_name
        enh_path = ENH_DIR / img_name
        
        with Image.open(orig_path) as img:
            # 1. AUTO-ROTATE (Fixes Vol 2-7 landscape scans)
            # If width > height, rotate 90 degrees. expand=True prevents cropping.
            if img.width > img.height:
                img = img.rotate(90, expand=True)
            
            # 2. GRAYSCALE (Removes yellowing and paper noise)
            img = ImageOps.grayscale(img)
            
            # 3. CONTRAST BOOST (Makes faded text stand out)
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(1.8) # 1.8 is aggressive for readability
            
            # 4. SHARPEN (Clarifies letter edges)
            sharp = ImageEnhance.Sharpness(img)
            img = sharp.enhance(2.0)

            img.save(enh_path, "JPEG", quality=95)

        # 5. ADD TO PPTX
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Center horizontally on a 13.33" slide
        left = Inches(3.6) 
        top = Inches(0.25)
        # Fix height to 7" to leave small margins top/bottom
        slide.shapes.add_picture(str(enh_path), left, top, height=Inches(7.0))

    prs.save(OUTPUT_PPTX)
    print(f"Done! Enhanced archive saved to: {OUTPUT_PPTX}")

if __name__ == "__main__":
    enhance_scans()
