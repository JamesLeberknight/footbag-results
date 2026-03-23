import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image  # You might need to: pip install Pillow

ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = ROOT / "out" / "scans"
OUTPUT_PPTX = ROOT / "out" / "Footbag_Archive_Scans.pptx"

def create_pptx():
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    images = sorted([f for f in os.listdir(IMG_DIR) if f.lower().endswith(('.jpg', '.jpeg'))])

    for img_name in images:
        img_path = IMG_DIR / img_name
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # --- AUTO-ROTATE LOGIC ---
        with Image.open(img_path) as img:
            width, height = img.size
        
        # If it's a horizontal scan (landscape), rotate it!
        # Note: python-pptx doesn't rotate the file, it rotates the SHAPE on the slide.
        left = Inches(3.0) # Center it a bit more
        top = Inches(0.5)
        
        pic = slide.shapes.add_picture(str(img_path), left, top, height=Inches(6.5))
        
        if width > height:
            pic.rotation = 90  # Flips it to portrait orientation

    prs.save(OUTPUT_PPTX)
    print(f"Done! Created {len(images)} slides.")

if __name__ == "__main__":
    create_pptx()
