#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

VALID_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff")


def main():
    img_dir = Path("out/manual_crops_prepped")
    output_ppt = Path("out/review_slides.pptx")

    images = sorted([p for p in img_dir.iterdir() if p.suffix.lower() in VALID_EXTS])

    if not images:
        print("No images found.")
        return

    prs = Presentation()

    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)

        # Add image (we'll scale it)
        pic = slide.shapes.add_picture(str(img_path), 0, 0)

        # Scale to fit slide
        img_width = pic.width
        img_height = pic.height

        scale_w = slide_width / img_width
        scale_h = slide_height / img_height
        scale = min(scale_w, scale_h)

        pic.width = int(img_width * scale)
        pic.height = int(img_height * scale)

        # Center image
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = int((slide_height - pic.height) / 2)

        print(f"Added: {img_path.name}")

    prs.save(output_ppt)
    print(f"\nSaved PowerPoint: {output_ppt}")


if __name__ == "__main__":
    main()
