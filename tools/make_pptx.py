import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, reversed_compare

# --- CONFIG ---
ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = ROOT / "out" / "scans"
OUTPUT_PPTX = ROOT / "out" / "Footbag_Archive_Scans.pptx"

def create_pptx():
    prs = Presentation()
    
    # Use standard 16:9 or 4:3 (default)
    # To set 16:9: prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    # Get all jpg/jpeg files and sort them alphabetically
    valid_exts = ('.jpg', '.jpeg', '.JPG', '.JPEG')
    images = sorted([f for f in os.listdir(IMG_DIR) if f.endswith(valid_exts)])

    if not images:
        print(f"No images found in {IMG_DIR}")
        return

    print(f"Found {len(images)} images. Generating PowerPoint...")

    for img_name in images:
        img_path = IMG_DIR / img_name
        
        # Add a blank slide layout (6 is usually the index for a blank slide)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add the image. We set a high width so it scales to fit the page 
        # but keeps the aspect ratio.
        try:
            # Center the image roughly
            left = Inches(0.5)
            top = Inches(0.5)
            # Constrain height to 6.5 inches to leave room for margins
            pic = slide.shapes.add_picture(str(img_path), left, top, height=Inches(6.5))
        except Exception as e:
            print(f"Skipping {img_name} due to error: {e}")

    prs.save(OUTPUT_PPTX)
    print(f"Done! Saved to {OUTPUT_PPTX}")

if __name__ == "__main__":
    create_pptx()
