from pathlib import Path
from io import BytesIO

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches

IMAGE_DIR = Path("out/images")
OUTPUT_PPTX = Path("raw_images_for_manual_crop_rotation.pptx")

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MARGIN_IN = 0.3

# File suffixes we will scan for. Actual format support is checked with PIL.
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp", ".mpo"}

# Formats python-pptx can reliably embed
PPT_COMPATIBLE_FORMATS = {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}


def get_image_files(image_dir: Path):
    files = [p for p in image_dir.iterdir() if p.is_file() and p.suffix.lower() in IMAGE_EXTS]
    return sorted(files, key=lambda p: p.name.lower())


def open_image_safely(img_path: Path):
    try:
        img = Image.open(img_path)
        img.load()
        return img
    except UnidentifiedImageError as e:
        raise ValueError(f"Cannot identify image file: {img_path}") from e


def get_image_size(img_path: Path):
    with Image.open(img_path) as img:
        img.load()
        return img.size


def fit_image_to_slide(img_path: Path, slide_w_in: float, slide_h_in: float, margin_in: float):
    px_w, px_h = get_image_size(img_path)
    img_aspect = px_w / px_h
    max_w = slide_w_in - 2 * margin_in
    max_h = slide_h_in - 2 * margin_in
    slide_aspect = max_w / max_h

    if img_aspect > slide_aspect:
        width_in = max_w
        height_in = width_in / img_aspect
    else:
        height_in = max_h
        width_in = height_in * img_aspect

    return width_in, height_in


def add_title_textbox(slide, text: str, slide_width_in: float):
    left = Inches(0.2)
    top = Inches(0.05)
    width = Inches(slide_width_in - 0.4)
    height = Inches(0.3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = text


def image_to_ppt_source(img_path: Path):
    """
    Return either:
    - a filesystem path string for directly supported formats, or
    - a BytesIO PNG stream for unsupported formats like MPO/WebP/etc.
    """
    img = open_image_safely(img_path)
    fmt = (img.format or "").upper()

    if fmt in PPT_COMPATIBLE_FORMATS:
        return str(img_path)

    # Convert unsupported formats to PNG in memory
    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGB")

    bio = BytesIO()
    img.save(bio, format="PNG")
    bio.seek(0)
    return bio


def main():
    if not IMAGE_DIR.exists():
        raise FileNotFoundError(f"Image directory not found: {IMAGE_DIR}")

    image_files = get_image_files(IMAGE_DIR)
    if not image_files:
        raise FileNotFoundError(f"No supported image files found in: {IMAGE_DIR}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    blank_layout = prs.slide_layouts[6]

    added = 0
    skipped = []

    for img_path in image_files:
        try:
            width_in, height_in = fit_image_to_slide(
                img_path, SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN, MARGIN_IN
            )
            ppt_source = image_to_ppt_source(img_path)

            slide = prs.slides.add_slide(blank_layout)
            add_title_textbox(slide, img_path.name, SLIDE_WIDTH_IN)

            left_in = (SLIDE_WIDTH_IN - width_in) / 2
            top_in = (SLIDE_HEIGHT_IN - height_in) / 2 + 0.15

            slide.shapes.add_picture(
                ppt_source,
                Inches(left_in),
                Inches(top_in),
                width=Inches(width_in),
                height=Inches(height_in),
            )
            added += 1

        except Exception as e:
            skipped.append((img_path.name, str(e)))

    prs.save(OUTPUT_PPTX)

    print(f"Created: {OUTPUT_PPTX}")
    print(f"Slides added: {added}")
    if skipped:
        print("\nSkipped files:")
        for name, err in skipped:
            print(f"  - {name}: {err}")


if __name__ == "__main__":
    main()
